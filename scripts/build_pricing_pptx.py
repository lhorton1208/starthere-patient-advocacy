#!/usr/bin/env python3
"""Build a StartHere Patient Advocacy pricing PowerPoint from site pricing data."""

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Brand colors from pricing.html
TEAL = RGBColor(0x1A, 0x6E, 0x6E)
TEAL_LIGHT = RGBColor(0xE1, 0xF5, 0xEE)
ORANGE = RGBColor(0xD4, 0x83, 0x3A)
WARM_BADGE = RGBColor(0xFA, 0xEE, 0xDA)
TEXT = RGBColor(0x2C, 0x2C, 0x2C)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG = RGBColor(0xF7, 0xF5, 0xF2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, text, size=18, bold=False, color=TEXT, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size, bold, color)
    return box


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rounded_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shape, fill)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def style_chart(chart, series_colors):
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(11)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = TEXT
    plot.data_labels.number_format = '"$"#,##0'
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    for i, series in enumerate(chart.series):
        color = series_colors[i % len(series_colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.color.rgb = MUTED
    chart.value_axis.tick_labels.number_format = '"$"#,##0'
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.category_axis.tick_labels.font.color.rgb = TEXT
    chart.category_axis.tick_labels.font.bold = True


def blank_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    # soft background
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    fill_shape(bg, SOFT_BG)
    # move to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    # teal accent bar at top
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.12))
    fill_shape(bar, TEAL)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---- Slide 1: Title ----
    slide = blank_slide(prs)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.5),
                "START HERE PATIENT ADVOCACY", size=14, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1),
                "Pricing Overview", size=44, bold=True, color=TEXT)
    add_textbox(slide, Inches(0.8), Inches(3.9), Inches(10), Inches(0.8),
                "Advocacy services, priced for the support you need.\nAll services are private pay.", size=18, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
                "Free 30-minute initial consultation  ·  No insurance billing  ·  Flexible engagement options",
                size=14, color=TEAL)

    # ---- Slide 2: Three tiers overview cards ----
    slide = blank_slide(prs)
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.4),
                "Service structures", size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.55),
                "Three ways to engage", size=28, bold=True, color=TEXT)

    tiers = [
        {
            "badge": "Pay as you go",
            "name": "Hourly advocacy",
            "price": "$95–$125/hr",
            "sub": "Billed in 30-minute increments",
            "desc": "One-time needs — a single appointment, a specific question, or a short-term issue.",
            "features": [
                "Appointment accompaniment",
                "Insurance / billing questions",
                "No long-term commitment",
            ],
            "featured": False,
        },
        {
            "badge": "Most chosen",
            "name": "Monthly care navigation",
            "price": "$600–$1,800/mo",
            "sub": "6–15 hours included monthly",
            "desc": "Ongoing coordination across providers, appointments, and paperwork.",
            "features": [
                "Dedicated advocate & phone line",
                "Scheduling & coordination",
                "Family updates",
                "Priority response (24 hrs)",
            ],
            "featured": True,
        },
        {
            "badge": "Project-based",
            "name": "Fixed-fee packages",
            "price": "$300–$2,500",
            "sub": "Per engagement, scoped upfront",
            "desc": "Defined situations with a clear beginning and end — hospital stay or procedure.",
            "features": [
                "Flat fee, no surprises",
                "Acute & episodic events",
                "Scoped before start",
            ],
            "featured": False,
        },
    ]

    card_w = Inches(3.85)
    card_h = Inches(5.3)
    gap = Inches(0.25)
    start_x = Inches(0.7)
    top = Inches(1.45)

    for i, tier in enumerate(tiers):
        left = start_x + i * (card_w + gap)
        card = add_rounded_rect(slide, left, top, card_w, card_h, WHITE)
        if tier["featured"]:
            card.line.color.rgb = TEAL
            card.line.width = Pt(2.5)

        badge_fill = WARM_BADGE if tier["featured"] else TEAL_LIGHT
        badge_color = RGBColor(0x63, 0x38, 0x06) if tier["featured"] else RGBColor(0x08, 0x50, 0x41)
        badge = add_rounded_rect(slide, left + Inches(0.25), top + Inches(0.25), Inches(1.6), Inches(0.32), badge_fill)
        try:
            badge.adjustments[0] = 0.5
        except Exception:
            pass
        add_textbox(slide, left + Inches(0.3), top + Inches(0.27), Inches(1.5), Inches(0.3),
                    tier["badge"], size=10, bold=True, color=badge_color, align=PP_ALIGN.CENTER)

        add_textbox(slide, left + Inches(0.25), top + Inches(0.7), Inches(3.3), Inches(0.35),
                    tier["name"], size=16, bold=True, color=TEXT)
        add_textbox(slide, left + Inches(0.25), top + Inches(1.1), Inches(3.3), Inches(0.7),
                    tier["desc"], size=12, color=MUTED)
        add_textbox(slide, left + Inches(0.25), top + Inches(1.85), Inches(3.3), Inches(0.45),
                    tier["price"], size=24, bold=True, color=TEAL)
        add_textbox(slide, left + Inches(0.25), top + Inches(2.3), Inches(3.3), Inches(0.3),
                    tier["sub"], size=11, color=MUTED)

        feat_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(2.75), Inches(3.3), Inches(2.2))
        tf = feat_box.text_frame
        tf.word_wrap = True
        first = True
        for feat in tier["features"]:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(8)
            run = p.add_run()
            set_run(run, f"✓  {feat}", size=12, color=TEXT)

    # ---- Slide 3: Tier price range chart ----
    slide = blank_slide(prs)
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.4),
                "Pricing comparison", size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.55),
                "Tier price ranges at a glance", size=28, bold=True, color=TEXT)
    add_textbox(slide, Inches(0.7), Inches(1.25), Inches(12), Inches(0.35),
                "Low and high ends of each engagement option (USD)", size=14, color=MUTED)

    chart_data = CategoryChartData()
    chart_data.categories = ["Hourly advocacy\n($/hr)", "Monthly care\nnavigation ($/mo)", "Fixed-fee\npackages"]
    chart_data.add_series("Low end", (95, 600, 300))
    chart_data.add_series("High end", (125, 1800, 2500))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.2), Inches(1.8), Inches(10.8), Inches(5.0),
        chart_data,
    ).chart
    style_chart(chart, [TEAL, ORANGE])

    # ---- Slide 4: Fixed-fee packages chart ----
    slide = blank_slide(prs)
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.4),
                "Fixed-fee packages", size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.55),
                "Common situations, scoped as a flat fee", size=28, bold=True, color=TEXT)

    packages = [
        ("ER Visit", 500, 1500),
        ("Inpatient Stay", 600, 2500),
        ("Hospital discharge", 750, 2000),
        ("Outpatient procedure", 350, 900),
        ("After-encounter follow-up", 300, 750),
    ]

    chart_data = CategoryChartData()
    chart_data.categories = [p[0] for p in packages]
    chart_data.add_series("Low end", [p[1] for p in packages])
    chart_data.add_series("High end", [p[2] for p in packages])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.4),
        chart_data,
    ).chart
    style_chart(chart, [TEAL, ORANGE])

    # ---- Slide 5: Package details table-style ----
    slide = blank_slide(prs)
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.4),
                "Fixed-fee packages", size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.5),
                "What’s included in each package", size=28, bold=True, color=TEXT)

    pkg_details = [
        ("ER Visit", "$500–$1,500",
         "On-site or remote support during an ER visit — communicating with staff, tracking decisions, keeping family informed."),
        ("Inpatient Stay", "$600–$2,500",
         "Ongoing presence throughout a hospital stay — daily updates, care team communication, decision support."),
        ("Hospital discharge", "$750–$2,000",
         "Safe transition home or to rehab, including follow-up care setup and home health referrals."),
        ("Outpatient procedure", "$350–$900",
         "Accompaniment and coordination before, during, and after a scheduled outpatient procedure."),
        ("After-encounter follow-up", "$300–$750",
         "Review discharge instructions, schedule follow-ups, confirm nothing falls through the cracks."),
    ]

    # header row
    y = Inches(1.4)
    header = add_rounded_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.45), TEAL)
    add_textbox(slide, Inches(0.9), y + Inches(0.08), Inches(3.2), Inches(0.35),
                "Package", size=13, bold=True, color=WHITE)
    add_textbox(slide, Inches(4.2), y + Inches(0.08), Inches(2.2), Inches(0.35),
                "Price range", size=13, bold=True, color=WHITE)
    add_textbox(slide, Inches(6.6), y + Inches(0.08), Inches(5.7), Inches(0.35),
                "Description", size=13, bold=True, color=WHITE)

    row_h = Inches(0.95)
    for i, (name, price, desc) in enumerate(pkg_details):
        ry = y + Inches(0.5) + i * row_h
        bg = WHITE if i % 2 == 0 else TEAL_LIGHT
        add_rounded_rect(slide, Inches(0.7), ry, Inches(11.9), row_h - Inches(0.08), bg)
        add_textbox(slide, Inches(0.9), ry + Inches(0.25), Inches(3.2), Inches(0.4),
                    name, size=14, bold=True, color=TEXT)
        add_textbox(slide, Inches(4.2), ry + Inches(0.25), Inches(2.2), Inches(0.4),
                    price, size=14, bold=True, color=TEAL)
        add_textbox(slide, Inches(6.6), ry + Inches(0.12), Inches(5.7), Inches(0.75),
                    desc, size=12, color=MUTED)

    # ---- Slide 6: Closing / notes ----
    slide = blank_slide(prs)
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.4),
                "Getting started", size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.55),
                "Key points for clients", size=28, bold=True, color=TEXT)

    notes = [
        ("Free initial consultation",
         "Every relationship starts with a complimentary 30-minute conversation to understand the situation and recommend the right level of support — no obligation."),
        ("Private pay only",
         "StartHere Patient Advocacy is not a medical provider and does not bill insurance or Medicare directly. All services are private pay."),
        ("Flexible structures",
         "Choose hourly for one-time needs, a monthly retainer for ongoing navigation, or a fixed-fee package for a defined care event."),
    ]

    for i, (title, body) in enumerate(notes):
        ny = Inches(1.55) + i * Inches(1.55)
        add_rounded_rect(slide, Inches(0.7), ny, Inches(11.9), Inches(1.35), WHITE)
        accent = slide.shapes.add_shape(1, Inches(0.7), ny, Inches(0.12), Inches(1.35))
        fill_shape(accent, TEAL if i != 1 else ORANGE)
        add_textbox(slide, Inches(1.2), ny + Inches(0.25), Inches(10.8), Inches(0.35),
                    title, size=18, bold=True, color=TEXT)
        add_textbox(slide, Inches(1.2), ny + Inches(0.65), Inches(10.8), Inches(0.55),
                    body, size=14, color=MUTED)

    out = Path(__file__).resolve().parents[1] / "presentations" / "StartHere-Pricing.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out}")
    return out


if __name__ == "__main__":
    build()
